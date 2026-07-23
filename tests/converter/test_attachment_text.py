"""Tests for src.converter.attachment_text (PR-2b).

覆盖:
    - extract_text 派发逻辑 (扩展名 → 对应 extractor)
    - 4 个 extractor: txt / docx / xlsx / pptx (程序化生成 fixture)
    - PDF: pypdf 装着但跳过真实抽取测试 (生成有内容 PDF 太复杂),
      只测 missing file + 空 PDF 错误路径
    - _truncate 边界 (ATTACHMENT_TEXT_MAX_BYTES)
    - unsupported 扩展 / file not found 路径
"""

from __future__ import annotations

from pathlib import Path

import pytest

from src.converter.attachment_text import (
    ATTACHMENT_TEXT_MAX_BYTES,
    ExtractResult,
    _truncate,
    extract_text,
)


class TestDispatch:
    def test_missing_file(self, tmp_path: Path):
        r = extract_text(tmp_path / "nope.pdf")
        assert r.status == 'failed'
        assert 'file not found' in (r.error_message or '')
        assert r.text == ''
        assert r.extractor == 'none'

    def test_unsupported_extension(self, tmp_path: Path):
        f = tmp_path / "demo.bin"
        f.write_bytes(b"\x00\x01\x02\x03")
        r = extract_text(f)
        assert r.status == 'unsupported'
        assert r.text == ''
        assert 'unsupported extension' in (r.error_message or '')

    def test_extension_from_filename_param(self, tmp_path: Path):
        """file_path 没扩展时, filename arg 提供 ext."""
        f = tmp_path / "noext"
        f.write_text("hello world", encoding='utf-8')
        # 用 filename hint .txt
        r = extract_text(f, filename="demo.txt")
        assert r.status == 'extracted'
        assert 'hello world' in r.text


class TestPlaintext:
    def test_plain_txt_extraction(self, tmp_path: Path):
        f = tmp_path / "notes.txt"
        f.write_text("first line\nsecond line", encoding='utf-8')
        r = extract_text(f)
        assert r.status == 'extracted'
        assert r.extractor == 'plaintext'
        assert 'first line' in r.text
        assert 'second line' in r.text
        assert not r.truncated

    def test_csv_treated_as_plaintext(self, tmp_path: Path):
        f = tmp_path / "data.csv"
        f.write_text("a,b,c\n1,2,3\n", encoding='utf-8')
        r = extract_text(f)
        assert r.status == 'extracted'
        assert r.extractor == 'plaintext'
        assert 'a,b,c' in r.text

    def test_markdown_treated_as_plaintext(self, tmp_path: Path):
        f = tmp_path / "spec.md"
        f.write_text("# Heading\n\nbody text here.\n", encoding='utf-8')
        r = extract_text(f)
        assert r.status == 'extracted'
        assert r.extractor == 'plaintext'

    def test_empty_plaintext_failed(self, tmp_path: Path):
        f = tmp_path / "blank.txt"
        f.write_text("", encoding='utf-8')
        r = extract_text(f)
        assert r.status == 'failed'

    def test_plaintext_truncation(self, tmp_path: Path):
        f = tmp_path / "big.txt"
        # ATTACHMENT_TEXT_MAX_BYTES + 1024 over cap
        oversize = "x" * (ATTACHMENT_TEXT_MAX_BYTES + 1024)
        f.write_text(oversize, encoding='utf-8')
        r = extract_text(f)
        assert r.status == 'extracted'
        assert r.truncated is True
        assert len(r.text.encode('utf-8')) <= ATTACHMENT_TEXT_MAX_BYTES


class TestDocx:
    def test_docx_paragraphs_and_table(self, tmp_path: Path):
        import docx as python_docx
        doc = python_docx.Document()
        doc.add_paragraph("一段中文")
        doc.add_paragraph("English paragraph")
        tbl = doc.add_table(rows=2, cols=2)
        tbl.cell(0, 0).text = "Col A"
        tbl.cell(0, 1).text = "Col B"
        tbl.cell(1, 0).text = "v1"
        tbl.cell(1, 1).text = "v2"
        f = tmp_path / "demo.docx"
        doc.save(str(f))

        r = extract_text(f)
        assert r.status == 'extracted'
        assert r.extractor == 'docx'
        assert '一段中文' in r.text
        assert 'English paragraph' in r.text
        assert 'Col A' in r.text and 'v2' in r.text
        # table 行用 ' | ' 拼接
        assert ' | ' in r.text

    def test_empty_docx_failed(self, tmp_path: Path):
        import docx as python_docx
        doc = python_docx.Document()
        f = tmp_path / "empty.docx"
        doc.save(str(f))

        r = extract_text(f)
        assert r.status == 'failed'
        assert 'empty' in (r.error_message or '').lower()


class TestPptx:
    def test_pptx_extract_slide_text(self, tmp_path: Path):
        import pptx as python_pptx
        from pptx.util import Inches
        pres = python_pptx.Presentation()
        layout = pres.slide_layouts[5]  # title-only layout
        slide = pres.slides.add_slide(layout)
        slide.shapes.title.text = "本周计划"
        # add text shape
        txbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
        txbox.text_frame.text = "redis cluster scaling"
        f = tmp_path / "deck.pptx"
        pres.save(str(f))

        r = extract_text(f)
        assert r.status == 'extracted'
        assert r.extractor == 'pptx'
        assert 'Slide 1' in r.text
        assert '本周计划' in r.text
        assert 'redis cluster scaling' in r.text

    def test_pptx_multiple_slides(self, tmp_path: Path):
        import pptx as python_pptx
        pres = python_pptx.Presentation()
        for i in range(3):
            slide = pres.slides.add_slide(pres.slide_layouts[5])
            slide.shapes.title.text = f"Slide title {i+1}"
        f = tmp_path / "multi.pptx"
        pres.save(str(f))

        r = extract_text(f)
        assert r.status == 'extracted'
        # 3 个 slide marker
        for i in range(1, 4):
            assert f'## Slide {i}' in r.text


class TestXlsx:
    def test_xlsx_extract_sheets(self, tmp_path: Path):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Plan"
        ws.append(["Project", "Owner", "Deadline"])
        ws.append(["redis migration", "alice", "2026-06-01"])
        ws.append(["产品评审", "bob", "2026-05-30"])
        # 第二个 sheet
        ws2 = wb.create_sheet("Notes")
        ws2.append(["misc note here"])
        f = tmp_path / "plan.xlsx"
        wb.save(str(f))

        r = extract_text(f)
        assert r.status == 'extracted'
        assert r.extractor == 'xlsx'
        assert 'Plan' in r.text
        assert 'Notes' in r.text
        assert 'redis migration' in r.text
        assert '产品评审' in r.text
        assert 'misc note here' in r.text
        # 行 用 ' | ' 拼接
        assert ' | ' in r.text

    def test_empty_xlsx_failed(self, tmp_path: Path):
        from openpyxl import Workbook
        wb = Workbook()
        # active sheet 默认空 — 删了它再加一个空 sheet 保证 calamine 看到的全空
        f = tmp_path / "empty.xlsx"
        wb.save(str(f))

        r = extract_text(f)
        # 空 sheet 在 calamine 路径下 to_python 返 [] → status='failed'
        assert r.status == 'failed'
        assert 'empty xlsx' in (r.error_message or '').lower()


class TestPdf:
    def test_pdf_missing_pypdf_or_empty(self, tmp_path: Path):
        """合法的空 PDF (no extractable text) → 'failed'.
        生成一个 minimal valid PDF 用 pypdf.PdfWriter.add_blank_page.
        """
        try:
            from pypdf import PdfWriter
        except ImportError:
            pytest.skip("pypdf not installed")

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        f = tmp_path / "blank.pdf"
        with open(f, 'wb') as fh:
            writer.write(fh)

        r = extract_text(f)
        # blank page extract_text 返 '' → status='failed'
        assert r.status == 'failed'
        assert r.extractor == 'pypdf'
        assert 'no extractable text' in (r.error_message or '')

    def test_pdf_not_actually_pdf_format(self, tmp_path: Path):
        """非 PDF 文件改成 .pdf 扩展名 → pypdf 抛 except → status='failed'."""
        f = tmp_path / "fake.pdf"
        f.write_bytes(b"not a pdf, just text bytes")
        r = extract_text(f)
        assert r.status == 'failed'
        # error_message 由 pypdf parse 错误填充
        assert r.error_message is not None


class TestTruncate:
    def test_truncate_short_text(self):
        text, truncated = _truncate("hello")
        assert text == "hello"
        assert truncated is False

    def test_truncate_long_text(self):
        long_text = "a" * (ATTACHMENT_TEXT_MAX_BYTES + 100)
        text, truncated = _truncate(long_text)
        assert truncated is True
        assert len(text.encode('utf-8')) <= ATTACHMENT_TEXT_MAX_BYTES

    def test_truncate_empty(self):
        text, truncated = _truncate("")
        assert text == ""
        assert truncated is False

    def test_truncate_multibyte_safe(self):
        # 中文 3 字节/字; 给一个紧贴 cap 的长度
        cjk = "中" * 100_000   # 300000 bytes > 256 KB
        text, truncated = _truncate(cjk)
        assert truncated is True
        # decoded 不爆 (errors='replace' 兜底 multibyte 切断)
        assert len(text.encode('utf-8')) <= ATTACHMENT_TEXT_MAX_BYTES


class TestExtractResultDataclass:
    def test_default_fields(self):
        r = ExtractResult(text='x', extractor='pypdf', status='extracted')
        assert r.error_message is None
        assert r.truncated is False


class TestImageOcr:
    """图片附件 OCR 派发 (批次4 PR-G) —— mock vision_ocr, 不真跑 Vision。"""

    def _png(self, tmp_path: Path, name: str = 'scan.png', data: bytes = b'\x89PNGfake') -> Path:
        f = tmp_path / name
        f.write_bytes(data)
        return f

    def test_image_extracted(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_image_file', lambda p: '合同条款 redis timeout')
        f = self._png(tmp_path)
        r = extract_text(f)
        assert r.status == 'extracted'
        assert r.extractor == 'vision_ocr'
        assert '合同条款' in r.text

    def test_flag_off_falls_through_unsupported(self, tmp_path: Path, monkeypatch):
        """flag off → 与现状逐字节一致 (unsupported extension, extractor none)。"""
        from src.converter import attachment_text
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: False)
        f = self._png(tmp_path)
        r = extract_text(f)
        assert r.status == 'unsupported'
        assert r.extractor == 'none'
        assert 'unsupported extension' in (r.error_message or '')

    def test_ocr_unavailable_falls_through_unsupported(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: False)
        f = self._png(tmp_path, 'scan.jpg')
        r = extract_text(f)
        assert r.status == 'unsupported'
        assert r.extractor == 'none'
        assert 'unsupported extension' in (r.error_message or '')

    def test_ocr_runtime_none_falls_through_unsupported(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_image_file', lambda p: None)
        f = self._png(tmp_path)
        r = extract_text(f)
        assert r.status == 'unsupported'
        assert r.extractor == 'none'

    def test_empty_ocr_result_unsupported(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_image_file', lambda p: '   ')
        f = self._png(tmp_path)
        r = extract_text(f)
        assert r.status == 'unsupported'
        assert r.extractor == 'vision_ocr'
        assert 'no text' in (r.error_message or '').lower()

    def test_oversized_image_unsupported(self, tmp_path: Path, monkeypatch):
        """>OCR_IMAGE_MAX_BYTES → unsupported + 'too large' 说明 (护栏)。"""
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'OCR_IMAGE_MAX_BYTES', 4)
        f = self._png(tmp_path, data=b'12345678')  # 8 bytes > 4
        r = extract_text(f)
        assert r.status == 'unsupported'
        assert r.extractor == 'vision_ocr'
        assert 'too large' in (r.error_message or '')

    def test_all_image_extensions_dispatch(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_image_file', lambda p: 'text here')
        for ext in attachment_text.OCR_IMAGE_EXTENSIONS:
            f = tmp_path / f'img{ext}'
            f.write_bytes(b'data')
            r = extract_text(f)
            assert r.status == 'extracted', ext
            assert r.extractor == 'vision_ocr', ext


class TestPdfOcrCascade:
    """无文本层 PDF 级联 OCR (批次4 PR-G) —— mock ocr_pdf_file。"""

    def _blank_pdf(self, tmp_path: Path) -> Path:
        from pypdf import PdfWriter
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        f = tmp_path / 'scan.pdf'
        with open(f, 'wb') as fh:
            writer.write(fh)
        return f

    def test_cascade_extracted(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_pdf_file', lambda p: ('扫描件正文 scanned', False))
        r = extract_text(self._blank_pdf(tmp_path))
        assert r.status == 'extracted'
        assert r.extractor == 'pdf_ocr'
        assert '扫描件正文' in r.text
        assert r.truncated is False

    def test_cascade_truncated(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_pdf_file', lambda p: ('page text', True))
        r = extract_text(self._blank_pdf(tmp_path))
        assert r.status == 'extracted'
        assert r.extractor == 'pdf_ocr'
        assert r.truncated is True

    def test_flag_off_keeps_failed_message(self, tmp_path: Path, monkeypatch):
        """flag off → failed 文案逐字节不变 (无 OCR 级联)。"""
        from src.converter import attachment_text
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: False)
        r = extract_text(self._blank_pdf(tmp_path))
        assert r.status == 'failed'
        assert r.extractor == 'pypdf'
        assert r.error_message == 'no extractable text (image-only PDF / 扫描件?)'

    def test_ocr_unavailable_keeps_failed_message(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: False)
        r = extract_text(self._blank_pdf(tmp_path))
        assert r.status == 'failed'
        assert r.error_message == 'no extractable text (image-only PDF / 扫描件?)'

    def test_ocr_returns_none_keeps_failed_message(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_pdf_file', lambda p: None)
        r = extract_text(self._blank_pdf(tmp_path))
        assert r.status == 'failed'
        assert r.error_message == 'no extractable text (image-only PDF / 扫描件?)'

    def test_ocr_empty_text_keeps_failed_message(self, tmp_path: Path, monkeypatch):
        from src.converter import attachment_text, vision_ocr
        monkeypatch.setattr(attachment_text, '_ocr_enabled', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_available', lambda: True)
        monkeypatch.setattr(vision_ocr, 'ocr_pdf_file', lambda p: ('   ', False))
        r = extract_text(self._blank_pdf(tmp_path))
        assert r.status == 'failed'
        assert r.error_message == 'no extractable text (image-only PDF / 扫描件?)'


class TestLegacyOffice:
    """老 Office (.doc/.ppt/.xls) soffice 桥 (批次4 PR-H) —— mock _run_soffice_convert。

    mock 造一个真实的转换产物 (docx/pptx/xlsx), 让下游原生 extractor 真跑, 验证
    「soffice 桥 → 现成 extractor」端到端接线 + extractor='soffice_bridge' 归一。
    """

    def _make_input(self, tmp_path: Path, name: str) -> Path:
        f = tmp_path / name
        f.write_bytes(b'\xd0\xcf\x11\xe0legacy-office-binary')  # OLE2-ish magic
        return f

    def _fake_convert(self, produce):
        """返回 fake _run_soffice_convert: produce(out_path, format) 造真实产物 → True。"""
        def _fake(input_path, output_dir, format='pdf', timeout=120):
            out = Path(output_dir) / f'{Path(input_path).stem}.{format}'
            produce(out, format)
            return True
        return _fake

    def test_doc_bridge_extracts(self, tmp_path: Path, monkeypatch):
        from src.converter import office_converter

        def produce(out, fmt):
            import docx as python_docx
            assert fmt == 'docx'
            doc = python_docx.Document()
            doc.add_paragraph('合同条款 from legacy doc')
            doc.save(str(out))

        monkeypatch.setattr(
            office_converter, '_run_soffice_convert', self._fake_convert(produce)
        )
        r = extract_text(self._make_input(tmp_path, 'contract.doc'))
        assert r.status == 'extracted'
        assert r.extractor == 'soffice_bridge'
        assert '合同条款' in r.text

    def test_ppt_bridge_extracts(self, tmp_path: Path, monkeypatch):
        from src.converter import office_converter

        def produce(out, fmt):
            import pptx as python_pptx
            assert fmt == 'pptx'
            pres = python_pptx.Presentation()
            slide = pres.slides.add_slide(pres.slide_layouts[5])
            slide.shapes.title.text = '季度汇报'
            pres.save(str(out))

        monkeypatch.setattr(
            office_converter, '_run_soffice_convert', self._fake_convert(produce)
        )
        r = extract_text(self._make_input(tmp_path, 'deck.ppt'))
        assert r.status == 'extracted'
        assert r.extractor == 'soffice_bridge'
        assert 'Slide 1' in r.text
        assert '季度汇报' in r.text

    def test_xls_bridge_extracts(self, tmp_path: Path, monkeypatch):
        from src.converter import office_converter

        def produce(out, fmt):
            from openpyxl import Workbook
            assert fmt == 'xlsx'
            wb = Workbook()
            ws = wb.active
            ws.append(['产品', '数量'])
            ws.append(['redis migration', 3])
            wb.save(str(out))

        monkeypatch.setattr(
            office_converter, '_run_soffice_convert', self._fake_convert(produce)
        )
        r = extract_text(self._make_input(tmp_path, 'plan.xls'))
        assert r.status == 'extracted'
        assert r.extractor == 'soffice_bridge'
        assert 'redis migration' in r.text
        assert '产品' in r.text

    def test_soffice_missing_or_failed_unsupported(self, tmp_path: Path, monkeypatch):
        """soffice 缺失 / 转换失败 (_run_soffice_convert → False) → unsupported, 不崩。"""
        from src.converter import office_converter
        monkeypatch.setattr(
            office_converter, '_run_soffice_convert', lambda *a, **k: False
        )
        r = extract_text(self._make_input(tmp_path, 'old.doc'))
        assert r.status == 'unsupported'
        assert r.extractor == 'soffice_bridge'
        assert 'LibreOffice' in (r.error_message or '')

    def test_soffice_true_but_no_output_unsupported(self, tmp_path: Path, monkeypatch):
        """_run_soffice_convert 返 True 却没产物 → unsupported (不误报 extracted)。"""
        from src.converter import office_converter
        monkeypatch.setattr(
            office_converter, '_run_soffice_convert', lambda *a, **k: True
        )
        r = extract_text(self._make_input(tmp_path, 'old.xls'))
        assert r.status == 'unsupported'
        assert r.extractor == 'soffice_bridge'
        assert 'produced no' in (r.error_message or '')

    def test_all_legacy_extensions_dispatch(self, tmp_path: Path, monkeypatch):
        """三个老格式扩展名都进 soffice 桥 (不落回 'unsupported extension' 现状)。"""
        from src.converter import attachment_text, office_converter
        monkeypatch.setattr(
            office_converter, '_run_soffice_convert', lambda *a, **k: False
        )
        for ext in attachment_text.LEGACY_OFFICE_EXTENSIONS:
            f = self._make_input(tmp_path, f'legacy{ext}')
            r = extract_text(f)
            assert r.extractor == 'soffice_bridge', ext
            assert r.status == 'unsupported', ext

    def test_real_soffice_doc_smoke(self, tmp_path: Path):
        """本机装了 LibreOffice → 真转一个最小 .doc 端到端; 没装 → skip。"""
        from src.converter.office_converter import _find_soffice, _run_soffice_convert
        if _find_soffice() is None:
            pytest.skip('LibreOffice/soffice not installed')
        import docx as python_docx
        doc = python_docx.Document()
        doc.add_paragraph('real soffice bridge smoke 合同 redis')
        docx_path = tmp_path / 'seed.docx'
        doc.save(str(docx_path))
        # docx → .doc via soffice, 再对那个 .doc 跑 extract_text 走桥回来。
        ok = _run_soffice_convert(str(docx_path), str(tmp_path), format='doc')
        doc_path = tmp_path / 'seed.doc'
        if not ok or not doc_path.exists():
            pytest.skip('soffice could not produce a .doc for smoke')
        r = extract_text(doc_path)
        assert r.status == 'extracted'
        assert r.extractor == 'soffice_bridge'
        assert '合同' in r.text or 'smoke' in r.text
