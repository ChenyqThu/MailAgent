"""Attachment text extraction — PDF / docx / pptx / xlsx → plaintext (PR-2b).

为 FTS5 索引和 LLM context 提取附件可搜索文本.

支持:
    PDF (.pdf):  pypdf.PdfReader, 按页抽 text + join '\\n\\n'; 无文本层(扫描件)
                 级联 Vision OCR (见 vision_ocr.py, flag MAILAGENT_ATTACHMENT_OCR_ENABLED)
    DOCX (.docx): python-docx Document, paragraphs + tables → markdown
    PPTX (.pptx): python-pptx Presentation, 按 slide 抽 title + body shapes
    XLSX (.xlsx): python-calamine 拼 markdown table (复用 office_converter
                  现成 Rust 引擎; fallback 到 pandas + openpyxl)
    图片 (png/jpg/…): macOS Vision OCR 识别中英文本 (批次4 PR-G, 同 flag 门控)
    老 Office (.doc/.ppt/.xls): soffice 桥转成 docx/pptx/xlsx 后复用现成 extractor
                  (批次4 PR-H, extractor='soffice_bridge'; soffice 缺失 → unsupported)
    TXT / MD / CSV: 直接 read_text() (utf-8 with errors=replace)
    其他: unsupported (zip/二进制等不支持)

不做的事:
    - 表格保 formatting: xlsx 折叠 csv-like 文本, 不还原 cell merge / color
    - 二进制压缩包 (zip): 历史邮件可能跳过

设计:
    extract_text(file_path, content_type=None, filename=None) → ExtractResult
    单步同步, caller 决定是否 async / fire-and-forget / thread pool.

输出文本会被截到 ATTACHMENT_TEXT_MAX_BYTES (默认 256 KB utf-8), 防 PDF 巨型
导致 FTS5 索引膨胀; 截断时 ExtractResult.truncated=True.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from loguru import logger


# Extensions
_PDF_EXTENSIONS = {'.pdf'}
_DOCX_EXTENSIONS = {'.docx'}
_PPTX_EXTENSIONS = {'.pptx'}
_XLSX_EXTENSIONS = {'.xlsx'}
_PLAINTEXT_EXTENSIONS = {'.txt', '.md', '.csv', '.log'}

# 图片附件 OCR 派发集（批次4 PR-G）。**单源**：存量 requeue CLI（PR-H）import
# 同一常量圈选待重跑行，禁止第二份硬编码。
OCR_IMAGE_EXTENSIONS = {
    '.png', '.jpg', '.jpeg', '.gif', '.heic', '.webp', '.tiff', '.bmp',
}

# 老 Office 二进制格式派发集（批次4 PR-H）。soffice 桥转成对应新格式后复用现成
# extractor。**单源**：存量 requeue CLI import 同一常量圈选待重跑行，禁止第二份硬编码。
LEGACY_OFFICE_EXTENSIONS = {'.doc', '.ppt', '.xls'}

# 老格式 → soffice --convert-to 目标新格式（转出后喂对应新格式 extractor）。
_LEGACY_OFFICE_TARGET_FORMAT = {
    '.doc': 'docx',
    '.ppt': 'pptx',
    '.xls': 'xlsx',
}

# 256 KB utf-8 bytes — FTS5 单 doc 上限的一个友好 cap.
# 大 PDF 抽完整 text 可能 >1 MB, 直接索引会让 FTS5 体积爆炸 + bm25 评分
# 不准 (高频 term 在长文档过度加权). 256 KB ≈ 50-100 页 text-only PDF
# 已经覆盖 95%+ 邮件附件场景.
ATTACHMENT_TEXT_MAX_BYTES = 256 * 1024


@dataclass
class ExtractResult:
    """单次 extraction 的结果."""
    text: str                                    # 提取出来的文本; 失败 / unsupported → ''
    extractor: str                               # 'pypdf' / 'docx' / 'pptx' / 'xlsx' / 'plaintext' / 'none'
    status: str                                  # 'extracted' | 'failed' | 'unsupported'
    error_message: Optional[str] = None          # 失败时的错误描述
    truncated: bool = False                      # 是否被 ATTACHMENT_TEXT_MAX_BYTES 截断


def extract_text(
    file_path: Path | str,
    content_type: Optional[str] = None,
    filename: Optional[str] = None,
) -> ExtractResult:
    """根据扩展名派发到对应 extractor.

    Args:
        file_path: 附件落地的本地路径
        content_type: MIME (e.g. 'application/pdf'), 当前仅作 hint (ext 优先)
        filename: 原文件名, ext 判断 fallback (file_path 没扩展名时用)

    Returns:
        ExtractResult, text 不超过 ATTACHMENT_TEXT_MAX_BYTES utf-8 bytes.
    """
    p = Path(file_path)
    if not p.exists():
        return ExtractResult(
            text='', extractor='none', status='failed',
            error_message=f'file not found: {p}',
        )

    # 扩展名: 优先 file_path 本身, fallback 原 filename (因为
    # AttachmentStore 落盘时可能用 sanitized filename, 但通常带 ext)
    ext = (p.suffix or Path(filename or '').suffix or '').lower()

    try:
        if ext in _PDF_EXTENSIONS:
            return _extract_pdf(p)
        if ext in _DOCX_EXTENSIONS:
            return _extract_docx(p)
        if ext in _PPTX_EXTENSIONS:
            return _extract_pptx(p)
        if ext in _XLSX_EXTENSIONS:
            return _extract_xlsx(p)
        if ext in _PLAINTEXT_EXTENSIONS:
            return _extract_plaintext(p)
        if ext in OCR_IMAGE_EXTENSIONS:
            ocr = _extract_image_ocr(p)
            if ocr is not None:
                return ocr
            # OCR 关闭 / 不可用 → 落回下方 unsupported（与现状逐字节一致）
        if ext in LEGACY_OFFICE_EXTENSIONS:
            return _extract_legacy_office(p, ext)
        return ExtractResult(
            text='', extractor='none', status='unsupported',
            error_message=f'unsupported extension: {ext!r}',
        )
    except Exception as e:
        logger.warning(f'extract_text failed for {p.name} (ext={ext}): {e}')
        return ExtractResult(
            text='', extractor='none', status='failed',
            error_message=str(e),
        )


def _ocr_enabled() -> bool:
    """OCR 总开关（MAILAGENT_ATTACHMENT_OCR_ENABLED，默认 true）。

    读 config 单例（与其余 config 消费一致，不直读 os.environ）。
    """
    from src.config import config
    return bool(config.mailagent_attachment_ocr_enabled)


def _extract_image_ocr(path: Path) -> Optional[ExtractResult]:
    """图片附件 → Vision OCR（批次4 PR-G）。

    Returns:
        - 有文本 → ExtractResult(extractor='vision_ocr', status='extracted')
        - OCR 跑通但无文字 / 超大图片 → ExtractResult(status='unsupported')
        - flag off / OCR 不可用 → None（调用方落回 'unsupported extension' 现状语义）
    """
    if not _ocr_enabled():
        return None
    from src.converter import vision_ocr
    if not vision_ocr.ocr_available():
        return None

    size = path.stat().st_size
    if size > vision_ocr.OCR_IMAGE_MAX_BYTES:
        return ExtractResult(
            text='', extractor='vision_ocr', status='unsupported',
            error_message=(
                f'image too large for OCR: {size} bytes '
                f'> {vision_ocr.OCR_IMAGE_MAX_BYTES}'
            ),
        )

    ocr_text = vision_ocr.ocr_image_file(path)
    if ocr_text is None:
        # OCR 运行期异常 → 落回现状 unsupported（返回 None 让上层出统一文案）
        return None
    text, truncated = _truncate(ocr_text)
    if not text.strip():
        return ExtractResult(
            text='', extractor='vision_ocr', status='unsupported',
            error_message='OCR produced no text (image has no readable content?)',
        )
    return ExtractResult(
        text=text, extractor='vision_ocr', status='extracted',
        truncated=truncated,
    )


def _truncate(text: str) -> tuple[str, bool]:
    """utf-8 字节截断到 ATTACHMENT_TEXT_MAX_BYTES, character 边界友好.

    用 errors='ignore' 在切断处丢弃不完整 multi-byte char (而不是 'replace'
    引入 U+FFFD 替换符 — 它本身 3 字节 utf-8, 会让结果 over cap).
    """
    if not text:
        return text, False
    encoded = text.encode('utf-8', errors='replace')
    if len(encoded) <= ATTACHMENT_TEXT_MAX_BYTES:
        return text, False
    truncated = encoded[:ATTACHMENT_TEXT_MAX_BYTES].decode('utf-8', errors='ignore')
    return truncated, True


def _extract_pdf(path: Path) -> ExtractResult:
    """pypdf.PdfReader 按页抽 text. 图片 PDF / 扫描件抽不出 → 'failed'."""
    try:
        from pypdf import PdfReader
    except ImportError as e:
        return ExtractResult(text='', extractor='pypdf', status='failed',
                             error_message=f'pypdf not installed: {e}')

    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            t = page.extract_text() or ''
        except Exception as e:
            logger.debug(f'pypdf page {i} extract err in {path.name}: {e}')
            continue
        t = t.strip()
        if t:
            parts.append(t)

    raw = '\n\n'.join(parts)
    text, truncated = _truncate(raw)
    if not text:
        # 无文本层（扫描件 / image-only PDF）→ 级联 Vision OCR（批次4 PR-G）。
        # OCR 关闭 / 不可用 / 渲染失败 / 无文字 → None，落回下方 failed 现状文案（逐字节不变）。
        ocr = _extract_pdf_ocr(path)
        if ocr is not None:
            return ocr
        return ExtractResult(
            text='', extractor='pypdf', status='failed',
            error_message='no extractable text (image-only PDF / 扫描件?)',
        )
    return ExtractResult(
        text=text, extractor='pypdf', status='extracted',
        truncated=truncated,
    )


def _extract_pdf_ocr(path: Path) -> Optional[ExtractResult]:
    """无文本层 PDF → Quartz 逐页渲染 + Vision OCR（批次4 PR-G）。

    Returns:
        - 有文本 → ExtractResult(extractor='pdf_ocr', status='extracted')；
          页数超上限则 truncated=True
        - flag off / OCR 不可用 / 渲染失败 / OCR 无文字 → None
          （调用方维持 'no extractable text' failed 现状文案，逐字节不变）
    """
    if not _ocr_enabled():
        return None
    from src.converter import vision_ocr
    if not vision_ocr.ocr_available():
        return None

    out = vision_ocr.ocr_pdf_file(path)
    if out is None:
        return None
    ocr_text, page_truncated = out
    text, byte_truncated = _truncate(ocr_text)
    if not text.strip():
        return None
    return ExtractResult(
        text=text, extractor='pdf_ocr', status='extracted',
        truncated=byte_truncated or page_truncated,
    )


def _extract_legacy_office(path: Path, ext: str) -> ExtractResult:
    """老 Office 二进制 (.doc/.ppt/.xls) → soffice 桥转新格式 → 复用现成 extractor（批次4 PR-H）。

    ``office_converter._run_soffice_convert(format=...)`` 转成 docx/pptx/xlsx（产物落
    ``tempfile.TemporaryDirectory``，抽完出作用域即删），再喂对应 ``_extract_docx`` /
    ``_extract_pptx`` / ``_extract_xlsx``。抽取标识统一 ``soffice_bridge``（区别原生 extractor）。

    soffice 缺失 / 转换失败 / 无产物 → ``unsupported``（graceful，现状 skip 语义；
    ``office_converter._run_soffice_convert`` 自带 try/except，找不到 soffice 返回 False
    不抛异常）。转换成功但转出文档为空 → 沿用底层 extractor 的 ``failed`` 语义。
    """
    import tempfile

    from src.converter import office_converter

    target_format = _LEGACY_OFFICE_TARGET_FORMAT[ext]
    with tempfile.TemporaryDirectory(prefix='legacy_office_') as tmpdir:
        ok = office_converter._run_soffice_convert(
            str(path), tmpdir, format=target_format,
        )
        if not ok:
            return ExtractResult(
                text='', extractor='soffice_bridge', status='unsupported',
                error_message=(
                    f'legacy office conversion unavailable/failed for {ext} '
                    f'(requires LibreOffice/soffice)'
                ),
            )
        converted = Path(tmpdir) / f'{path.stem}.{target_format}'
        if not converted.exists() or converted.stat().st_size == 0:
            return ExtractResult(
                text='', extractor='soffice_bridge', status='unsupported',
                error_message=f'soffice produced no {target_format} output for {path.name}',
            )
        # TemporaryDirectory 作用域内读产物（extractor 一次性把 text 读进内存）。
        if target_format == 'docx':
            result = _extract_docx(converted)
        elif target_format == 'pptx':
            result = _extract_pptx(converted)
        else:
            result = _extract_xlsx(converted)

    # extractor 统一标 soffice_bridge；status / text / truncated / error 沿用底层结果。
    return ExtractResult(
        text=result.text, extractor='soffice_bridge', status=result.status,
        error_message=result.error_message, truncated=result.truncated,
    )


def _extract_docx(path: Path) -> ExtractResult:
    """python-docx: paragraphs + tables → markdown 拼接."""
    try:
        import docx as python_docx
    except ImportError as e:
        return ExtractResult(text='', extractor='docx', status='failed',
                             error_message=f'python-docx not installed: {e}')

    doc = python_docx.Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [
                cell.text.strip().replace('\n', ' ').replace('|', '\\|')
                for cell in row.cells
            ]
            line = ' | '.join(cells)
            if line.strip(' |'):
                parts.append(line)

    raw = '\n'.join(parts)
    text, truncated = _truncate(raw)
    if not text:
        return ExtractResult(
            text='', extractor='docx', status='failed',
            error_message='empty docx (no paragraphs or tables)',
        )
    return ExtractResult(
        text=text, extractor='docx', status='extracted',
        truncated=truncated,
    )


def _extract_pptx(path: Path) -> ExtractResult:
    """python-pptx: 按 slide → '## Slide N' heading + 各 shape text."""
    try:
        import pptx as python_pptx
    except ImportError as e:
        return ExtractResult(text='', extractor='pptx', status='failed',
                             error_message=f'python-pptx not installed: {e}')

    pres = python_pptx.Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(pres.slides, start=1):
        slide_parts = [f'## Slide {idx}']
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    slide_parts.append(t)
        if len(slide_parts) > 1:
            parts.append('\n'.join(slide_parts))

    raw = '\n\n'.join(parts)
    text, truncated = _truncate(raw)
    if not text:
        return ExtractResult(
            text='', extractor='pptx', status='failed',
            error_message='empty pptx (no slide with text)',
        )
    return ExtractResult(
        text=text, extractor='pptx', status='extracted',
        truncated=truncated,
    )


def _extract_xlsx(path: Path) -> ExtractResult:
    """xlsx via python-calamine — Rust 引擎 4-18× pandas/openpyxl 快.
    fallback 到 pandas+openpyxl 兼容 (但通常 calamine 都装上了, 见
    requirements.txt).
    """
    try:
        from python_calamine import CalamineWorkbook  # type: ignore[import-not-found]
    except ImportError:
        return _extract_xlsx_pandas(path)

    try:
        wb = CalamineWorkbook.from_path(str(path))
    except Exception as e:
        return ExtractResult(text='', extractor='xlsx', status='failed',
                             error_message=f'calamine open failed: {e}')

    parts: list[str] = []
    for sheet_name in wb.sheet_names:
        sheet = wb.get_sheet_by_name(sheet_name)
        try:
            rows = sheet.to_python()
        except Exception as e:
            logger.debug(f'calamine sheet {sheet_name!r} read err: {e}')
            continue
        if not rows:
            continue
        parts.append(f'## {sheet_name}')
        for row in rows:
            line = ' | '.join('' if c is None else str(c) for c in row)
            if line.strip(' |'):
                parts.append(line)

    raw = '\n\n'.join(parts)
    text, truncated = _truncate(raw)
    if not text:
        return ExtractResult(
            text='', extractor='xlsx', status='failed',
            error_message='empty xlsx (no sheet with data)',
        )
    return ExtractResult(
        text=text, extractor='xlsx', status='extracted',
        truncated=truncated,
    )


def _extract_xlsx_pandas(path: Path) -> ExtractResult:
    """pandas + openpyxl fallback — calamine 缺时用."""
    try:
        import pandas as pd
    except ImportError as e:
        return ExtractResult(text='', extractor='xlsx', status='failed',
                             error_message=f'pandas not installed: {e}')
    try:
        sheets = pd.read_excel(str(path), sheet_name=None)
    except Exception as e:
        return ExtractResult(text='', extractor='xlsx', status='failed',
                             error_message=f'pandas read_excel failed: {e}')

    parts: list[str] = []
    for name, df in sheets.items():
        if df.empty:
            continue
        parts.append(f'## {name}')
        parts.append(df.to_csv(index=False))

    raw = '\n\n'.join(parts)
    text, truncated = _truncate(raw)
    if not text:
        return ExtractResult(
            text='', extractor='xlsx', status='failed',
            error_message='empty xlsx (no sheet with data)',
        )
    return ExtractResult(
        text=text, extractor='xlsx', status='extracted',
        truncated=truncated,
    )


def _extract_plaintext(path: Path) -> ExtractResult:
    """直接 read_text utf-8 with errors=replace, 防 latin-1 / gbk 等编码烂."""
    try:
        raw = path.read_text(encoding='utf-8', errors='replace')
    except OSError as e:
        return ExtractResult(text='', extractor='plaintext', status='failed',
                             error_message=str(e))
    text, truncated = _truncate(raw)
    if not text.strip():
        return ExtractResult(text='', extractor='plaintext', status='failed',
                             error_message='empty file')
    return ExtractResult(
        text=text, extractor='plaintext', status='extracted',
        truncated=truncated,
    )
